#!/usr/bin/env python3
from __future__ import annotations

import argparse
from functools import lru_cache
import os
import re
from pathlib import Path
from typing import Iterable

import fitz  # pymupdf
from docx import Document
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation


NON_TRANSLATABLE_PATTERNS = [
    re.compile(r"https?://\S+"),
    re.compile(r"`[^`]+`"),
    re.compile(r"\b[A-Za-z]:\\[^\s]+"),
    re.compile(r"(?:^|\s)(?:-\w+|--\w[\w-]*)(?:\s|$)"),
    re.compile(r"\b(?:Claude|OpenAI|ChatGPT|Anthropic|API|SDK)\b"),
]

PDF_CJK_FONT = "china-t"
PDF_EMBEDDED_CJK_FONT_NAME = "zh_tw_cjk"
PDF_CJK_FONT_CANDIDATES = [
    r"C:\Windows\Fonts\msjh.ttc",
    r"C:\Windows\Fonts\msjhbd.ttc",
    r"C:\Windows\Fonts\NotoSansTC-VF.ttf",
]
LOCKED_TERM_NORMALIZERS = [
    (re.compile(r"C\s+l\s+a\s+u\s+d\s+e", re.IGNORECASE), "Claude"),
    (re.compile(r"O\s+p\s+e\s+n\s+A\s+I", re.IGNORECASE), "OpenAI"),
    (re.compile(r"C\s+h\s+a\s+t\s+G\s+P\s+T", re.IGNORECASE), "ChatGPT"),
]
LOCKED_TERM_SPLIT = re.compile(r"(Claude|OpenAI|ChatGPT|Anthropic|API|SDK)")
GUIDE_SOURCE_RE = re.compile(r"\b(complete\s+guide|guide|manual|handbook)\b", re.IGNORECASE)
GUIDE_TARGET_RE = re.compile(r"(完整指南|指南|手冊)")
ENGLISH_TERM_BEFORE_SKILL_RE = re.compile(
    r"\b([A-Z][A-Za-z0-9.+_-]*(?:\s+[A-Z][A-Za-z0-9.+_-]*){0,3})\s+技能\b"
)


def make_output_path(src: Path) -> Path:
    return src.with_name(f"{src.stem}.zh-tw{src.suffix}")


def should_translate(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if re.fullmatch(r"[\W\d_]+", stripped):
        return False
    return True


def looks_like_model_meta_response(text: str) -> bool:
    lowered = text.lower()
    meta_markers = [
        "please provide the text",
        "provide the text you want me to translate",
        "請提供您要翻譯",
        "請提供要翻譯",
        "無法翻譯",
    ]
    return any(marker in lowered for marker in meta_markers)


def mask_non_translatable(text: str) -> tuple[str, dict[str, str]]:
    mapping: dict[str, str] = {}
    masked = text
    idx = 0
    for pattern in NON_TRANSLATABLE_PATTERNS:
        matches = list(pattern.finditer(masked))
        for m in reversed(matches):
            token = m.group(0)
            key = f"__LOCK_{idx}__"
            idx += 1
            mapping[key] = token
            masked = masked[: m.start()] + key + masked[m.end() :]
    return masked, mapping


def unmask(text: str, mapping: dict[str, str]) -> str:
    out = text
    for key, token in mapping.items():
        out = out.replace(key, token)
    for pattern, replacement in LOCKED_TERM_NORMALIZERS:
        out = pattern.sub(replacement, out)
    return out


def normalize_translation(source_text: str, target_text: str) -> str:
    out = target_text
    source_lines = [line.strip() for line in source_text.splitlines() if line.strip()]
    target_lines = [line.strip() for line in out.splitlines() if line.strip()]
    if (
        len(source_lines) > 1
        and GUIDE_SOURCE_RE.search(source_lines[0])
        and target_lines
        and not GUIDE_TARGET_RE.search(target_lines[0])
    ):
        joined = " ".join(target_lines)
        match = GUIDE_TARGET_RE.search(joined)
        if match:
            guide_label = match.group(1)
            topic = GUIDE_TARGET_RE.sub("", joined, count=1).strip()
            topic = re.sub(r"的\s*$", "", topic).strip()
            out = f"{guide_label}\n{topic}" if topic else guide_label
    if re.search(r"\bfor\s+[A-Z][A-Za-z0-9.+_-]*", source_text, re.IGNORECASE):
        out = ENGLISH_TERM_BEFORE_SKILL_RE.sub(r"\1 的技能", out)
    return out


GUIDE_TARGET_RE = re.compile(r"(\u5b8c\u6574\u6307\u5357|\u6307\u5357|\u624b\u518a)")
ENGLISH_TERM_BEFORE_SKILL_RE = re.compile(
    r"\b([A-Z][A-Za-z0-9.+_-]*(?:\s+[A-Z][A-Za-z0-9.+_-]*){0,3})\s+\u6280\u80fd\b"
)


def normalize_translation(source_text: str, target_text: str) -> str:
    out = target_text
    source_lines = [line.strip() for line in source_text.splitlines() if line.strip()]
    target_lines = [line.strip() for line in out.splitlines() if line.strip()]
    if (
        len(source_lines) > 1
        and GUIDE_SOURCE_RE.search(source_lines[0])
        and target_lines
        and not GUIDE_TARGET_RE.search(target_lines[0])
    ):
        joined = " ".join(target_lines)
        match = GUIDE_TARGET_RE.search(joined)
        if match:
            guide_label = match.group(1)
            topic = GUIDE_TARGET_RE.sub("", joined, count=1).strip()
            topic = re.sub(r"\u7684\s*$", "", topic).strip()
            out = f"{guide_label}\n{topic}" if topic else guide_label
    if re.search(r"\bfor\s+[A-Z][A-Za-z0-9.+_-]*", source_text, re.IGNORECASE):
        out = ENGLISH_TERM_BEFORE_SKILL_RE.sub(
            lambda match: f"{match.group(1)} \u7684\u6280\u80fd", out
        )
    return out


def pdf_text_color(rgb: int) -> tuple[float, float, float]:
    r = ((rgb >> 16) & 255) / 255.0
    g = ((rgb >> 8) & 255) / 255.0
    b = (rgb & 255) / 255.0
    return (r, g, b)


@lru_cache(maxsize=1)
def pdf_cjk_fontfile() -> str | None:
    configured = os.getenv("PDF_CJK_FONT_PATH")
    candidates = [configured] if configured else []
    candidates.extend(PDF_CJK_FONT_CANDIDATES)
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def _average_page_color(
    page: fitz.Page, clips: list[fitz.Rect]
) -> tuple[float, float, float] | None:
    totals = [0, 0, 0]
    count = 0
    for clip in clips:
        clip &= page.rect
        if clip.is_empty or clip.width < 1 or clip.height < 1:
            continue
        pix = page.get_pixmap(clip=clip, alpha=False)
        samples = pix.samples
        stride = pix.n
        for idx in range(0, len(samples), stride):
            totals[0] += samples[idx]
            totals[1] += samples[idx + 1]
            totals[2] += samples[idx + 2]
            count += 1
    if count == 0:
        return None
    return tuple(v / count / 255.0 for v in totals)  # type: ignore[return-value]


def pdf_background_color(page: fitz.Page, rect: fitz.Rect) -> tuple[float, float, float]:
    pad = 3
    outer = rect + (-pad, -pad, pad, pad)
    clips = [
        fitz.Rect(outer.x0, outer.y0, outer.x1, rect.y0),
        fitz.Rect(outer.x0, rect.y1, outer.x1, outer.y1),
        fitz.Rect(outer.x0, rect.y0, rect.x0, rect.y1),
        fitz.Rect(rect.x1, rect.y0, outer.x1, rect.y1),
    ]
    try:
        sampled = _average_page_color(page, clips)
        if sampled is not None:
            return sampled
        sampled = _average_page_color(page, [rect])
        return sampled if sampled is not None else (1, 1, 1)
    except Exception:
        return (1, 1, 1)


def pdf_insert_fitted_text(
    page: fitz.Page,
    rect: fitz.Rect,
    text: str,
    font_size: float,
    color: tuple[float, float, float],
) -> bool:
    cjk_fontfile = pdf_cjk_fontfile() if any(ord(ch) > 127 for ch in text) else None
    fontname = PDF_EMBEDDED_CJK_FONT_NAME if cjk_fontfile else (
        PDF_CJK_FONT if any(ord(ch) > 127 for ch in text) else "helv"
    )
    size = max(4.0, min(48.0, float(font_size)))
    if cjk_fontfile:
        while size >= 4.0:
            result = page.insert_textbox(
                rect,
                text,
                fontsize=size,
                fontname=fontname,
                fontfile=cjk_fontfile,
                color=color,
                align=fitz.TEXT_ALIGN_LEFT,
            )
            if result >= 0:
                return True
            size -= 0.75
        return False

    if any(ord(ch) > 127 for ch in text) and LOCKED_TERM_SPLIT.search(text):
        while size >= 4.0:
            parts = [part for part in LOCKED_TERM_SPLIT.split(text) if part]
            widths = []
            for part in parts:
                part_font = "helv" if LOCKED_TERM_SPLIT.fullmatch(part) else PDF_CJK_FONT
                widths.append(fitz.get_text_length(part, fontname=part_font, fontsize=size))
            if sum(widths) <= rect.width:
                x = rect.x0
                y = min(rect.y1, rect.y0 + size)
                for part, width in zip(parts, widths):
                    part_font = "helv" if LOCKED_TERM_SPLIT.fullmatch(part) else PDF_CJK_FONT
                    page.insert_text(
                        (x, y),
                        part,
                        fontsize=size,
                        fontname=part_font,
                        color=color,
                    )
                    x += width
                return True
            size -= 0.75

        return False

    while size >= 4.0:
        result = page.insert_textbox(
            rect,
            text,
            fontsize=size,
            fontname=fontname,
            color=color,
            align=fitz.TEXT_ALIGN_LEFT,
        )
        if result >= 0:
            return True
        size -= 0.75

    return False


def _text_width(text: str, size: float) -> float:
    cjk_fontfile = pdf_cjk_fontfile() if any(ord(ch) > 127 for ch in text) else None
    if cjk_fontfile:
        return fitz.Font(fontfile=cjk_fontfile).text_length(text, fontsize=size)

    width = 0.0
    for part in [part for part in LOCKED_TERM_SPLIT.split(text) if part]:
        fontname = "helv" if LOCKED_TERM_SPLIT.fullmatch(part) else PDF_CJK_FONT
        width += fitz.get_text_length(part, fontname=fontname, fontsize=size)
    return width


def _wrap_text_for_rects(text: str, rects: list[tuple], font_size: float) -> list[str]:
    if len(rects) <= 1:
        return [text]
    explicit_lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(explicit_lines) > 1:
        return explicit_lines
    units = text.split()
    if len(units) <= 1:
        units = list(text)

    lines: list[str] = []
    cursor = ""
    rect_idx = 0
    for unit in units:
        candidate = f"{cursor} {unit}".strip() if text.split() else f"{cursor}{unit}"
        max_width = fitz.Rect(rects[min(rect_idx, len(rects) - 1)]).width
        if cursor and _text_width(candidate, font_size) > max_width and rect_idx < len(rects) - 1:
            lines.append(cursor)
            cursor = unit
            rect_idx += 1
        else:
            cursor = candidate
    if cursor:
        lines.append(cursor)
    return lines


def pdf_insert_text_in_line_rects(
    page: fitz.Page,
    line_rects: list[tuple],
    text: str,
    font_size: float,
    color: tuple[float, float, float],
) -> bool:
    lines = _wrap_text_for_rects(text, line_rects, font_size)
    fitted = True
    for idx, line in enumerate(lines[: len(line_rects)]):
        rect = fitz.Rect(line_rects[idx])
        if not pdf_insert_fitted_text(page, rect, line, font_size, color):
            fitted = False
    if len(lines) > len(line_rects):
        overflow = "".join(lines[len(line_rects) :])
        rect = fitz.Rect(line_rects[-1])
        if not pdf_insert_fitted_text(page, rect, overflow, max(4, font_size * 0.8), color):
            fitted = False
    return fitted


def pdf_can_fit_text_in_line_rects(
    line_rects: list[tuple], text: str, font_size: float, min_font_size: float = 4.0
) -> bool:
    lines = _wrap_text_for_rects(text, line_rects, font_size)
    if len(lines) > len(line_rects):
        return False
    for line, rect_tuple in zip(lines, line_rects):
        rect = fitz.Rect(rect_tuple)
        size = max(min_font_size, min(48.0, float(font_size)))
        while size >= min_font_size:
            if _text_width(line, size) <= rect.width:
                break
            size -= 0.75
        else:
            return False
    return True


def _looks_like_code_or_table(text: str, font_names: list[str]) -> bool:
    stripped = text.strip()
    if not stripped:
        return True
    if any("mono" in name.lower() or "code" in name.lower() for name in font_names):
        return True
    code_chars = sum(1 for ch in stripped if ch in "{}[]<>`|\\/_=+#:*;")
    alpha_chars = sum(1 for ch in stripped if ch.isalpha())
    if code_chars >= 6 and code_chars > alpha_chars * 0.35:
        return True
    if re.search(r"\b[a-zA-Z0-9_-]+\.(md|py|sh|yaml|yml|json|zip)\b", stripped):
        return True
    if re.search(r"^\s*[│├└─|+-]{2,}", stripped):
        return True
    return False


def _is_heading_block(block_lines: list[str], sizes: list[float]) -> bool:
    if not block_lines or not sizes:
        return False
    max_size = max(sizes)
    avg_size = sum(sizes) / len(sizes)
    joined = " ".join(block_lines)
    if GUIDE_SOURCE_RE.search(joined) and len(block_lines) <= 4:
        return True
    return len(block_lines) <= 3 and (max_size >= 22 or avg_size >= 18)


def pdf_line_records(page: fitz.Page) -> list[dict]:
    data = page.get_text("dict")
    records: list[dict] = []
    for block_idx, block in enumerate(data.get("blocks", [])):
        if block.get("type") != 0:
            continue
        block_lines: list[str] = []
        line_rects: list[tuple] = []
        line_records: list[dict] = []
        sizes: list[float] = []
        font_names: list[str] = []
        block_rect: fitz.Rect | None = None
        first_span: dict | None = None
        for line_idx, line in enumerate(block.get("lines", [])):
            spans = [s for s in line.get("spans", []) if s.get("text", "").strip()]
            if not spans:
                continue
            line_text = "".join(s.get("text", "") for s in spans).strip()
            if not line_text:
                continue
            line_rect = fitz.Rect(spans[0]["bbox"])
            for span in spans[1:]:
                line_rect |= fitz.Rect(span["bbox"])
            block_lines.append(line_text)
            line_rects.append(tuple(line_rect))
            block_rect = line_rect if block_rect is None else block_rect | line_rect
            first_span = first_span or spans[0]
            sizes.extend(float(span.get("size", 10)) for span in spans)
            font_names.extend(str(span.get("font", "")) for span in spans)
            line_records.append(
                {
                    "anchor_id": f"b{block_idx}.l{line_idx}",
                    "bbox": tuple(line_rect),
                    "line_rects": [tuple(line_rect)],
                    "text": line_text,
                    "size": spans[0].get("size", 10),
                    "color": spans[0].get("color", 0),
                    "skip": _looks_like_code_or_table(line_text, font_names),
                }
            )

        if block_lines and block_rect is not None and first_span is not None:
            if _looks_like_code_or_table("\n".join(block_lines), font_names):
                records.extend({**record, "skip": True} for record in line_records)
            elif _is_heading_block(block_lines, sizes):
                records.append(
                    {
                        "anchor_id": f"b{block_idx}",
                        "bbox": tuple(block_rect),
                        "line_rects": line_rects,
                        "text": "\n".join(block_lines),
                        "size": first_span.get("size", 10),
                        "color": first_span.get("color", 0),
                        "skip": False,
                    }
                )
            else:
                records.extend(line_records)
    return records


class Translator:
    def __init__(self, model: str) -> None:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError(
                "OPENAI_API_KEY is not set. Please set it before running translation."
            )
        self.client = OpenAI(api_key=api_key)
        self.model = model
        self.cache: dict[str, str] = {}

    def translate(self, text: str) -> str:
        if not should_translate(text):
            return text
        if text in self.cache:
            return self.cache[text]

        masked, mapping = mask_non_translatable(text)
        prompt = (
            "Translate the following text into Taiwan Traditional Chinese (zh-TW). "
            "Preserve meaning, keep placeholders like __LOCK_#__ unchanged, and do not add explanation. "
            "Prefer fluent Taiwan usage over literal word-by-word fragments, including natural possessive wording for phrases like 'X for Y'. "
            "For English title phrases shaped like 'object for product/person', prefer the zh-TW order 'product/person 的 object'. "
            "For heading blocks, preserve line roles: if a guide/manual/handbook label appears as its own source line, keep that document-type label as its own translated line instead of moving it to the end of another line. "
            "If the input contains line breaks, use the full context while preserving the same heading hierarchy and line-break intent where possible.\n\n"
            f"{masked}"
        )
        resp = self.client.responses.create(
            model=self.model,
            input=prompt,
            temperature=0,
        )
        translated = (resp.output_text or "").strip()
        if not translated or looks_like_model_meta_response(translated):
            translated = text
        translated = unmask(translated, mapping)
        translated = normalize_translation(text, translated)
        self.cache[text] = translated
        return translated


def localize_markdown(src: Path, out: Path, tr: Translator) -> None:
    lines = src.read_text(encoding="utf-8").splitlines()
    in_code = False
    output: list[str] = []
    for line in lines:
        if line.strip().startswith("```"):
            in_code = not in_code
            output.append(line)
            continue
        if in_code or line.strip().startswith("![]("):
            output.append(line)
            continue
        output.append(tr.translate(line))
    out.write_text("\n".join(output) + "\n", encoding="utf-8")


def _translate_docx_runs(runs: Iterable, tr: Translator) -> None:
    for run in runs:
        if should_translate(run.text):
            run.text = tr.translate(run.text)


def localize_docx(src: Path, out: Path, tr: Translator) -> None:
    doc = Document(str(src))
    for p in doc.paragraphs:
        _translate_docx_runs(p.runs, tr)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _translate_docx_runs(p.runs, tr)
    doc.save(str(out))


def localize_pptx(src: Path, out: Path, tr: Translator) -> None:
    prs = Presentation(str(src))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if should_translate(run.text):
                            run.text = tr.translate(run.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                if should_translate(run.text):
                                    run.text = tr.translate(run.text)
    prs.save(str(out))


def localize_pdf(src: Path, out: Path, tr: Translator, max_pages: int | None = None) -> None:
    doc = fitz.open(str(src))
    total_pages = len(doc)
    pages_to_process = total_pages if max_pages is None else min(max_pages, total_pages)

    for page_idx in range(pages_to_process):
        page = doc[page_idx]
        print(f"Processing page {page_idx + 1}/{pages_to_process}...", flush=True)
        spans_to_replace: list[dict] = []
        for line_record in pdf_line_records(page):
            if line_record.get("skip"):
                continue
            text = line_record["text"]
            if should_translate(text):
                translated = tr.translate(text)
                if translated != text and pdf_can_fit_text_in_line_rects(
                    line_record["line_rects"], translated, line_record["size"]
                ):
                    rect = fitz.Rect(line_record["bbox"])
                    spans_to_replace.append(
                        {
                            "bbox": line_record["bbox"],
                            "line_rects": line_record["line_rects"],
                            "text": translated,
                            "size": line_record["size"],
                            "color": line_record["color"],
                            "fill": pdf_background_color(page, rect),
                        }
                    )
        for s in spans_to_replace:
            rect = fitz.Rect(s["bbox"])
            page.add_redact_annot(rect, fill=s["fill"])
        if spans_to_replace:
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
        for s in spans_to_replace:
            pdf_insert_text_in_line_rects(
                page,
                s["line_rects"],
                s["text"],
                font_size=s["size"],
                color=pdf_text_color(s["color"]),
            )
    doc.save(str(out))
    doc.close()


def main() -> int:
    load_dotenv()
    parser = argparse.ArgumentParser(
        description="Localize PDF/Markdown/DOCX/PPTX into Taiwan Traditional Chinese."
    )
    parser.add_argument("input_file", help="Input file path")
    parser.add_argument(
        "--model",
        default=None,
        help="OpenAI model name. If omitted, use OPENAI_MODEL or fallback default.",
    )
    parser.add_argument(
        "--max-pages",
        type=int,
        default=None,
        help="Only process the first N pages for PDF. Ignored for non-PDF.",
    )
    args = parser.parse_args()

    src = Path(args.input_file).resolve()
    if not src.exists():
        raise FileNotFoundError(f"Input not found: {src}")

    out = make_output_path(src)
    model = args.model or os.getenv("OPENAI_MODEL", "gpt-4.1-mini")
    tr = Translator(model=model)

    ext = src.suffix.lower()
    if ext == ".md":
        localize_markdown(src, out, tr)
    elif ext == ".docx":
        localize_docx(src, out, tr)
    elif ext == ".pptx":
        localize_pptx(src, out, tr)
    elif ext == ".pdf":
        localize_pdf(src, out, tr, max_pages=args.max_pages)
    else:
        raise ValueError(f"Unsupported format: {ext}. Supported: .pdf .md .docx .pptx")

    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
